from __future__ import annotations

import argparse
import hashlib
import io
import json
import math
import mimetypes
import os
import re
import shutil
import subprocess
import sys
import importlib.util
import tempfile
import time
import uuid
import webbrowser
import zipfile
from collections import OrderedDict
from concurrent.futures import ThreadPoolExecutor, as_completed
from dataclasses import dataclass
from datetime import datetime
from html import escape
from http.server import SimpleHTTPRequestHandler, ThreadingHTTPServer
from pathlib import Path
from threading import Lock, Thread, Timer
from typing import Any
from urllib.parse import parse_qs, urlparse
import xml.etree.ElementTree as ET

from PIL import ExifTags, Image, ImageDraw, ImageFilter, ImageFont, ImageOps


APP_DIR = Path(__file__).resolve().parent
STATIC_DIR = APP_DIR / "static"
LOGO_DIR = APP_DIR / "logos"
LOGO_RULES_FILE = LOGO_DIR / "logo-rules.json"
LOCAL_PACKAGES_DIR = APP_DIR.parent / ".python-packages"
if LOCAL_PACKAGES_DIR.exists():
    sys.path.insert(0, str(LOCAL_PACKAGES_DIR))
IMAGE_EXTENSIONS = {".jpg", ".jpeg"}
LOGO_EXTENSIONS = {".png", ".jpg", ".jpeg", ".webp", ".bmp"}
SORT_MODES = {"dateAsc", "dateDesc", "nameAsc", "nameDesc"}
SESSIONS: dict[str, dict[str, Any]] = {}
EXPORT_JOBS: dict[str, dict[str, Any]] = {}
EXPORT_LOCK = Lock()
LOGO_CACHE: dict[str, tuple[tuple[int, int], Image.Image]] = {}
LOGO_CACHE_LOCK = Lock()
LOGO_RULES_CACHE: tuple[int, list[dict[str, Any]]] | None = None
LOGO_RULES_LOCK = Lock()
EXIF_CACHE_LOCK = Lock()
EXIF_CACHE_SCHEMA = 1
EXIF_CACHE_MAX_ENTRIES = 50000
PREVIEW_CACHE: OrderedDict[str, bytes] = OrderedDict()
PREVIEW_CACHE_LOCK = Lock()
PREVIEW_CACHE_BYTES = 0
MIB = 1024 * 1024
GIB = 1024 * MIB
PREVIEW_CACHE_DEFAULT_BYTES = 1024 * MIB
PREVIEW_CACHE_MIN_BYTES = 64 * MIB
PREVIEW_CACHE_MAX_BYTES = 4 * GIB
PREVIEW_CACHE_MIN_ITEMS = 256
PREVIEW_CACHE_MAX_ITEMS = 4096
PREVIEW_CACHE_ESTIMATED_ITEM_BYTES = 384 * 1024
RESAMPLE = getattr(getattr(Image, "Resampling", Image), "LANCZOS")
EXPORT_MAX_EDGE = 30000
EXPORT_MAX_PIXELS = 250_000_000

try:
    from pptx import Presentation
except Exception:
    Presentation = None

try:
    from lens_profiles import LENS_VALUE_TABLES, NIKON_LENS_IDS
except Exception:
    profile_spec = importlib.util.spec_from_file_location("lens_profiles", APP_DIR / "lens_profiles.py")
    if profile_spec and profile_spec.loader:
        lens_profiles = importlib.util.module_from_spec(profile_spec)
        profile_spec.loader.exec_module(lens_profiles)
        NIKON_LENS_IDS = getattr(lens_profiles, "NIKON_LENS_IDS", {})
        LENS_VALUE_TABLES = getattr(lens_profiles, "LENS_VALUE_TABLES", [])
    else:
        NIKON_LENS_IDS = {}
        LENS_VALUE_TABLES = []


@dataclass
class Photo:
    index: int
    path: str
    name: str
    size: int
    width: int
    height: int
    exif: dict[str, str]


class ExportCancelled(RuntimeError):
    pass


DEFAULT_SETTINGS = {
    "slideWidth": 1920,
    "slideHeight": 1080,
    "slideLongEdge": 1920,
    "slideAspectRatio": "16:9",
    "background": "#f6f4ef",
    "safeWidthPct": 95,
    "safeHeightPct": 87,
    "pageMarginXPct": 2.5,
    "pageMarginYPct": 1.85,
    "bannerWidthPct": 62.5,
    "bannerHeightPct": 7.407,
    "gapPct": 1.852,
    "bannerOpacity": 0.6,
    "bannerColor": "#000000",
    "infoFontPct": 2.037,
    "paramFontPct": 2.593,
    "brandFontPct": 5.556,
    "logoPath": "",
    "brandText": "",
    "bannerTextOverrides": {},
    "shadow": True,
    "quality": 92,
    "exportFormat": "jpeg",
    "exportScalePct": 100,
}


def list_images(folder: Path, recursive: bool) -> list[Path]:
    iterator = folder.rglob("*") if recursive else folder.iterdir()
    return sorted(
        [p for p in iterator if p.is_file() and p.suffix.lower() in IMAGE_EXTENSIONS],
        key=lambda p: natural_key(p.name),
    )


def natural_key(value: str) -> list[Any]:
    return [int(text) if text.isdigit() else text.casefold() for text in re.split(r"(\d+)", value)]


def normalize_sort_mode(value: Any) -> str:
    text = clean_text(value)
    return text if text in SORT_MODES else "dateAsc"


def find_exiftool(root: Path) -> str | None:
    candidates = [
        os.environ.get("EXIFTOOL_PATH"),
        str(root / "exiftool.exe"),
        str(APP_DIR.parent / "tools" / "exiftool" / "exiftool.exe"),
        shutil.which("exiftool"),
        shutil.which("exiftool.exe"),
    ]
    for candidate in candidates:
        if candidate and Path(candidate).exists():
            return candidate
    return None


def read_exif_batch(paths: list[Path], root: Path, recursive: bool) -> tuple[dict[str, dict[str, str]], str]:
    native_results: dict[str, dict[str, str]] = {}
    native_failures = 0
    cached_results, uncached_paths, exif_cache = cached_exif_results(paths)
    native_results.update(cached_results)
    for path, metadata, failed in read_native_exif_parallel(uncached_paths):
        native_results[str(path.resolve())] = metadata
        if failed:
            native_failures += 1
    if uncached_paths:
        remember_exif_results(exif_cache, native_results, uncached_paths)
    if any(has_useful_exif(metadata) for metadata in native_results.values()):
        source = "Native Python"
        if cached_results:
            source += " cache"
        if native_failures:
            source += f" ({native_failures} failed)"
        return native_results, source

    exiftool = find_exiftool(root)
    if exiftool:
        try:
            exiftool_results = read_with_exiftool(paths, exiftool, root, recursive)
            remember_exif_results(exif_cache, exiftool_results, paths)
            return exiftool_results, f"ExifTool ({exiftool})"
        except Exception:
            pass
    return native_results, "Native Python"


def cached_exif_results(paths: list[Path]) -> tuple[dict[str, dict[str, str]], list[Path], dict[str, Any]]:
    exif_cache = load_exif_cache()
    entries = exif_cache.setdefault("entries", {})
    cached: dict[str, dict[str, str]] = {}
    missing: list[Path] = []
    for path in paths:
        key, signature = exif_cache_key_and_signature(path)
        entry = entries.get(key)
        if isinstance(entry, dict) and entry.get("signature") == signature and isinstance(entry.get("exif"), dict):
            cached[str(path.resolve())] = {str(k): clean_text(v) for k, v in entry["exif"].items()}
        else:
            missing.append(path)
    return cached, missing, exif_cache


def remember_exif_results(exif_cache: dict[str, Any], results: dict[str, dict[str, str]], paths: list[Path]) -> None:
    entries = exif_cache.setdefault("entries", {})
    now = time.time()
    changed = False
    for path in paths:
        key, signature = exif_cache_key_and_signature(path)
        metadata = results.get(str(path.resolve()), results.get(str(path), {}))
        entries[key] = {
            "signature": signature,
            "exif": metadata,
            "used": now,
        }
        changed = True
    if changed:
        save_exif_cache(exif_cache)


def exif_cache_key_and_signature(path: Path) -> tuple[str, dict[str, Any]]:
    try:
        resolved = str(path.resolve())
    except OSError:
        resolved = str(path)
    if os.name == "nt":
        resolved = resolved.casefold()
    try:
        stat = path.stat()
        signature = {"mtime": stat.st_mtime_ns, "size": stat.st_size}
    except OSError:
        signature = {"mtime": 0, "size": 0}
    key = hashlib.sha256(resolved.encode("utf-8", errors="surrogatepass")).hexdigest()
    return key, signature


def load_exif_cache() -> dict[str, Any]:
    path = exif_cache_path()
    if path is None:
        return {"schema": EXIF_CACHE_SCHEMA, "entries": {}}
    with EXIF_CACHE_LOCK:
        try:
            data = json.loads(path.read_text(encoding="utf-8"))
            if data.get("schema") == EXIF_CACHE_SCHEMA and isinstance(data.get("entries"), dict):
                return data
        except Exception:
            pass
    return {"schema": EXIF_CACHE_SCHEMA, "entries": {}}


def save_exif_cache(exif_cache: dict[str, Any]) -> None:
    path = exif_cache_path()
    if path is None:
        return
    with EXIF_CACHE_LOCK:
        entries = exif_cache.get("entries")
        if not isinstance(entries, dict):
            return
        if len(entries) > EXIF_CACHE_MAX_ENTRIES:
            ordered = sorted(
                entries.items(),
                key=lambda item: float(item[1].get("used", 0)) if isinstance(item[1], dict) else 0,
                reverse=True,
            )
            exif_cache["entries"] = dict(ordered[:EXIF_CACHE_MAX_ENTRIES])
        path.parent.mkdir(parents=True, exist_ok=True)
        tmp = path.with_suffix(".tmp")
        tmp.write_text(json.dumps(exif_cache, ensure_ascii=False, separators=(",", ":")), encoding="utf-8")
        tmp.replace(path)


def exif_cache_path() -> Path | None:
    directory = app_cache_dir()
    return directory / "exif-cache.json" if directory else None


def app_cache_dir() -> Path | None:
    candidates: list[Path] = []
    override = os.environ.get("EXIF_BANNER_CACHE_DIR")
    if override:
        candidates.append(Path(override))
    if os.name == "nt":
        for value in (os.environ.get("LOCALAPPDATA"), os.environ.get("APPDATA")):
            if value:
                candidates.append(Path(value) / "EXIF-Banner")
    else:
        candidates.append(Path(os.environ.get("XDG_CACHE_HOME", Path.home() / ".cache")) / "EXIF-Banner")
    candidates.append(Path.home() / ".exif-banner")
    for path in candidates:
        try:
            path.mkdir(parents=True, exist_ok=True)
            return path
        except OSError:
            continue
    return None


def read_native_exif_parallel(paths: list[Path]) -> list[tuple[Path, dict[str, str], bool]]:
    if not paths:
        return []
    workers = background_worker_count(len(paths))
    if workers <= 1:
        return [read_native_exif_item(path) for path in paths]
    with ThreadPoolExecutor(max_workers=workers) as executor:
        return list(executor.map(read_native_exif_item, paths))


def read_native_exif_item(path: Path) -> tuple[Path, dict[str, str], bool]:
    try:
        return path, read_with_native_python(path), False
    except Exception:
        return path, {}, True


def background_worker_count(total: int) -> int:
    cpu_count = os.cpu_count() or 2
    workers = max(1, cpu_count - 1)
    return min(workers, max(1, total))


def has_useful_exif(metadata: dict[str, str]) -> bool:
    return any(metadata.get(key) for key in ("model", "lens", "focalLength", "fNumber", "exposureTime", "iso", "dateTime"))


def read_with_exiftool(paths: list[Path], exiftool: str, root: Path, recursive: bool) -> dict[str, dict[str, str]]:
    tags = [
        "-SourceFile",
        "-FileName",
        "-Make",
        "-Model",
        "-LensID",
        "-LensModel",
        "-ImageWidth",
        "-ImageHeight",
        "-FocalLength",
        "-FNumber",
        "-ExposureTime",
        "-ISO",
        "-DateTimeOriginal",
    ]
    try:
        command = [
            exiftool,
            "-j",
            "-charset",
            "filename=utf8",
            "-ext",
            "jpg",
            "-ext",
            "jpeg",
            *(["-r"] if recursive else []),
            *tags,
            ".",
        ]
        completed = subprocess.run(
            command,
            cwd=str(root),
            stdout=subprocess.PIPE,
            stderr=subprocess.PIPE,
            text=True,
            encoding="utf-8",
            errors="replace",
            check=True,
        )
        results = parse_exiftool_json(completed.stdout, root)
        if results:
            return results
    except Exception:
        pass
    return read_with_exiftool_paths(paths, exiftool, tags)


def read_with_exiftool_paths(paths: list[Path], exiftool: str, tags: list[str]) -> dict[str, dict[str, str]]:
    results: dict[str, dict[str, str]] = {}
    for chunk in chunked(paths, 35):
        command = [
            exiftool,
            "-j",
            "-charset",
            "filename=utf8",
            *tags,
            *[str(path) for path in chunk],
        ]
        completed = subprocess.run(
            command,
            stdout=subprocess.PIPE,
            stderr=subprocess.PIPE,
            text=True,
            encoding="utf-8",
            errors="replace",
            check=True,
        )
        results.update(parse_exiftool_json(completed.stdout, None))
    return results


def parse_exiftool_json(output: str, root: Path | None) -> dict[str, dict[str, str]]:
    results: dict[str, dict[str, str]] = {}
    for item in json.loads(output or "[]"):
        source_text = clean_text(item.get("SourceFile"))
        if not source_text:
            continue
        source = Path(source_text)
        if not source.is_absolute() and root is not None:
            source = root / source_text
        results[str(source.resolve())] = normalize_exif(item)
    return results


def chunked(values: list[Path], size: int) -> list[list[Path]]:
    return [values[index : index + size] for index in range(0, len(values), size)]


def read_with_native_python(path: Path) -> dict[str, str]:
    with Image.open(path) as image:
        raw: dict[str, Any] = {}
        raw["ImageWidth"], raw["ImageHeight"] = image.size
        exif = image.getexif()
        for key, value in exif.items():
            raw[ExifTags.TAGS.get(key, str(key))] = value
        for ifd_name in ("Exif", "GPSInfo"):
            try:
                ifd_id = getattr(ExifTags.IFD, ifd_name)
                ifd = exif.get_ifd(ifd_id)
                for key, value in ifd.items():
                    raw[ExifTags.TAGS.get(key, str(key))] = value
            except Exception:
                continue
    raw.update(read_xmp_metadata(path))
    return normalize_exif(raw)


def read_with_pillow(path: Path) -> dict[str, str]:
    return read_with_native_python(path)


XMP_HEADER = b"http://ns.adobe.com/xap/1.0/\x00"
XMP_FIELDS = {"Lens", "LensID", "LensInfo", "LensMake", "LensModel", "LensType", "Make", "Model"}


def read_xmp_metadata(path: Path) -> dict[str, str]:
    metadata: dict[str, str] = {}
    for payload in iter_jpeg_app1_segments(path):
        if payload.startswith(XMP_HEADER):
            metadata.update(parse_xmp_packet(payload[len(XMP_HEADER) :]))
    return metadata


def iter_jpeg_app1_segments(path: Path):
    with path.open("rb") as handle:
        if handle.read(2) != b"\xff\xd8":
            return
        while True:
            prefix = handle.read(1)
            if not prefix:
                return
            if prefix != b"\xff":
                continue
            marker_byte = handle.read(1)
            while marker_byte == b"\xff":
                marker_byte = handle.read(1)
            if not marker_byte:
                return
            marker = marker_byte[0]
            if marker in {0xD9, 0xDA}:
                return
            length_data = handle.read(2)
            if len(length_data) != 2:
                return
            length = int.from_bytes(length_data, "big")
            if length < 2:
                return
            payload = handle.read(length - 2)
            if marker == 0xE1:
                yield payload


def parse_xmp_packet(payload: bytes) -> dict[str, str]:
    text = payload.decode("utf-8", errors="replace").strip("\ufeff\x00\r\n ")
    metadata: dict[str, str] = {}
    try:
        root = ET.fromstring(text)
        for element in root.iter():
            for key, value in element.attrib.items():
                local = xml_local_name(key)
                if local in XMP_FIELDS:
                    metadata[f"XMP:{local}"] = clean_text(value)
            local = xml_local_name(element.tag)
            if local in XMP_FIELDS and element.text:
                metadata[f"XMP:{local}"] = clean_text(element.text)
    except ET.ParseError:
        for field in XMP_FIELDS:
            match = re.search(rf"(?:aux|exifEX|tiff):{field}\s*=\s*\"([^\"]+)\"", text)
            if match:
                metadata[f"XMP:{field}"] = clean_text(match.group(1))
    return metadata


def xml_local_name(name: str) -> str:
    if "}" in name:
        return name.rsplit("}", 1)[1]
    return name.rsplit(":", 1)[-1]


def normalize_exif(raw: dict[str, Any]) -> dict[str, str]:
    lens = resolve_lens_model(raw)
    return {
        "make": clean_text(pick(raw, "Make")),
        "model": clean_text(pick(raw, "Model")),
        "lens": clean_text(lens),
        "focalLength": format_focal(pick(raw, "FocalLength")),
        "fNumber": format_number(pick(raw, "FNumber")),
        "exposureTime": format_shutter(pick(raw, "ExposureTime")),
        "iso": clean_text(pick(raw, "ISO", "ISOSpeedRatings", "PhotographicSensitivity")),
        "dateTime": clean_text(pick(raw, "DateTimeOriginal", "CreateDate")),
        "imageWidth": clean_text(pick(raw, "ImageWidth", "ExifImageWidth")),
        "imageHeight": clean_text(pick(raw, "ImageHeight", "ExifImageHeight")),
    }


def resolve_lens_model(raw: dict[str, Any]) -> str:
    explicit_lens_id = clean_text(pick(raw, "Composite:LensID", "LensID"))
    if explicit_lens_id and not looks_like_numeric_id(explicit_lens_id):
        return explicit_lens_id

    make = clean_text(pick(raw, "Make", "XMP:Make"))
    lens_info = pick(raw, "XMP:LensInfo", "LensSpecification", "LensInfo")
    lens_text = pick(raw, "XMP:LensModel", "LensModel", "XMP:Lens", "Lens")
    if "nikon" in make.lower():
        lens_id = parse_lens_id_number(pick(raw, "XMP:LensID", "LensIDNumber", "LensID"))
        if lens_id is not None:
            resolved = resolve_nikon_lens(
                lens_id,
                lens_info,
                lens_text,
            )
            if resolved:
                return resolved
            return f"Nikon LensID {lens_id}"

    resolved = resolve_lens_from_value_tables(make, raw, lens_info, lens_text)
    if resolved:
        return resolved

    return clean_text(pick(raw, "LensModel", "XMP:LensModel", "XMP:Lens", "Lens", "LensInfo", "LensSpecification"))


def resolve_nikon_lens(lens_id: int, lens_info: Any, lens_text: Any) -> str:
    prefix = f"{lens_id:02X} "
    candidates = [(key, name) for key, name in NIKON_LENS_IDS.items() if key.startswith(prefix)]
    if not candidates:
        return ""
    unique_names = sorted({name for _, name in candidates})
    if len(unique_names) == 1:
        return unique_names[0]

    signature = lens_signature_from_info(lens_info) or lens_signature_from_text(clean_text(lens_text))
    if signature:
        matches = sorted({name for _, name in candidates if lens_signature_from_text(name) == signature})
        if len(matches) == 1:
            return matches[0]
        collapsed = collapse_lens_variant_names(matches)
        if collapsed:
            return collapsed
    return ""


def collapse_lens_variant_names(names: list[str]) -> str:
    if len(names) < 2:
        return ""
    ordered = sorted(set(names), key=lambda value: (len(value), value.casefold()))
    base = ordered[0]
    suffixes: list[str] = []
    for name in ordered[1:]:
        if not name.startswith(f"{base} "):
            return ""
        suffix = name[len(base) :].strip(" -/")
        if not suffix or not is_lens_variant_suffix(suffix):
            return ""
        suffixes.append(suffix)
    return f"{base} / {' / '.join(suffixes)}" if suffixes else ""


def is_lens_variant_suffix(value: str) -> bool:
    tokens = {token.upper() for token in re.findall(r"[A-Za-z0-9]+", value)}
    return bool(tokens) and tokens.issubset(
        {
            "APO",
            "II",
            "III",
            "IS",
            "OIS",
            "OS",
            "OSS",
            "VC",
            "VR",
        }
    )


BRAND_ALIASES = {
    "canon": ("canon",),
    "fujifilm": ("fuji", "fujifilm"),
    "minolta": ("minolta", "sony"),
    "nikon": ("nikon",),
    "olympus": ("olympus", "om digital"),
    "panasonic": ("panasonic", "leica"),
    "pentax": ("pentax", "ricoh"),
    "ricoh": ("ricoh", "pentax"),
    "samsung": ("samsung",),
    "sigma": ("sigma",),
    "sony": ("sony", "minolta"),
}


def resolve_lens_from_value_tables(make: str, raw: dict[str, Any], lens_info: Any, lens_text: Any) -> str:
    values = lens_lookup_values(raw)
    if not values:
        return ""

    candidates: set[str] = set()
    for table in LENS_VALUE_TABLES:
        if not lens_table_matches_make(table, make):
            continue
        table_values = table.get("values", {})
        if not isinstance(table_values, dict):
            continue
        for value in values:
            match = table_values.get(value)
            if isinstance(match, str) and is_likely_lens_model(match):
                candidates.add(match)

    if len(candidates) == 1:
        return next(iter(candidates))

    signature = lens_signature_from_info(lens_info) or lens_signature_from_text(clean_text(lens_text))
    if signature:
        signature_matches = sorted(name for name in candidates if lens_signature_from_text(name) == signature)
        if len(signature_matches) == 1:
            return signature_matches[0]
    return ""


def lens_lookup_values(raw: dict[str, Any]) -> list[str]:
    values: list[str] = []
    for key in ("XMP:LensID", "XMP:LensType", "LensID", "LensType", "LensType2", "LensType3", "RFLensType"):
        value = clean_text(raw.get(key))
        if not value:
            continue
        values.append(value)
        numeric = parse_lens_id_number(value)
        if numeric is not None:
            values.extend([str(numeric), f"{numeric:02X}", f"{numeric:04X}"])
    return sorted(set(values), key=len, reverse=True)


def lens_table_matches_make(table: dict[str, Any], make: str) -> bool:
    lowered_make = make.lower()
    table_text = f"{table.get('table', '')} {table.get('group1', '')} {table.get('group2', '')}".lower()
    if not lowered_make:
        return False
    for canonical, aliases in BRAND_ALIASES.items():
        if any(alias in lowered_make for alias in aliases):
            return canonical in table_text or any(alias in table_text for alias in aliases)
    return any(part and part in table_text for part in re.split(r"\W+", lowered_make))


def is_likely_lens_model(value: str) -> bool:
    lowered = value.lower()
    if any(word in lowered for word in ("unknown", "reserved", "n/a", "none")):
        return False
    return "lens" in lowered or "mm" in lowered or any(
        brand in lowered
        for brand in ("canon", "nikkor", "nikon", "sony", "zeiss", "sigma", "tamron", "tokina", "olympus", "zuiko", "pentax", "panasonic", "leica")
    )


def parse_lens_id_number(value: Any) -> int | None:
    text = clean_text(value)
    if not text:
        return None
    match = re.fullmatch(r"0x([0-9a-fA-F]+)", text)
    if match:
        return int(match.group(1), 16)
    match = re.fullmatch(r"\d+", text)
    if match:
        return int(text)
    return None


def looks_like_numeric_id(value: str) -> bool:
    return parse_lens_id_number(value) is not None


def lens_signature_from_info(value: Any) -> tuple[str, str, str, str] | None:
    values: list[float] = []
    if isinstance(value, (tuple, list)):
        for item in value[:4]:
            numeric = to_float(item)
            if numeric is not None:
                values.append(numeric)
    else:
        for part in re.findall(r"\d+(?:\.\d+)?(?:/\d+(?:\.\d+)?)?", clean_text(value)):
            numeric = to_float_ratio(part)
            if numeric is not None:
                values.append(numeric)
    if len(values) < 4:
        return None
    return tuple(canonical_number(value) for value in values[:4])  # type: ignore[return-value]


def lens_signature_from_text(value: str) -> tuple[str, str, str, str] | None:
    text = value.lower().replace("–", "-").replace("—", "-")
    focal_range = re.search(r"(\d+(?:\.\d+)?)\s*-\s*(\d+(?:\.\d+)?)\s*mm", text)
    if focal_range:
        min_focal, max_focal = focal_range.groups()
    else:
        focal = re.search(r"(\d+(?:\.\d+)?)\s*mm", text)
        if not focal:
            return None
        min_focal = max_focal = focal.group(1)

    aperture_range = re.search(r"f/?\s*(\d+(?:\.\d+)?)\s*-\s*(\d+(?:\.\d+)?)", text)
    if aperture_range:
        min_aperture, max_aperture = aperture_range.groups()
    else:
        aperture = re.search(r"f/?\s*(\d+(?:\.\d+)?)", text)
        if not aperture:
            return None
        min_aperture = max_aperture = aperture.group(1)
    return tuple(canonical_number(float(value)) for value in (min_focal, max_focal, min_aperture, max_aperture))  # type: ignore[return-value]


def to_float_ratio(value: str) -> float | None:
    if "/" in value:
        numerator, denominator = value.split("/", 1)
        try:
            return float(numerator) / float(denominator)
        except Exception:
            return None
    try:
        return float(value)
    except Exception:
        return None


def canonical_number(value: float | str) -> str:
    numeric = float(value)
    if numeric.is_integer():
        return str(int(numeric))
    return f"{numeric:.2f}".rstrip("0").rstrip(".")


def pick(raw: dict[str, Any], *keys: str) -> Any:
    for key in keys:
        value = raw.get(key)
        if value not in (None, ""):
            return value
    return ""


def clean_text(value: Any) -> str:
    if value is None:
        return ""
    if isinstance(value, bytes):
        value = value.decode("utf-8", errors="replace")
    return str(value).strip()


def to_float(value: Any) -> float | None:
    if value in (None, ""):
        return None
    if isinstance(value, tuple) and len(value) == 2 and value[1]:
        return float(value[0]) / float(value[1])
    try:
        return float(value)
    except Exception:
        try:
            return float(value.numerator) / float(value.denominator)
        except Exception:
            return None


def format_number(value: Any) -> str:
    numeric = to_float(value)
    if numeric is None:
        return clean_text(value)
    return f"{numeric:.1f}".rstrip("0").rstrip(".")


def format_focal(value: Any) -> str:
    text = clean_text(value)
    if text and "mm" in text.lower():
        return text
    numeric = to_float(value)
    if numeric is None:
        return text
    return f"{numeric:.1f} mm"


def format_shutter(value: Any) -> str:
    text = clean_text(value)
    if "/" in text:
        return text.removesuffix(" sec")
    numeric = to_float(value)
    if numeric is None:
        return text.removesuffix(" sec")
    if numeric and numeric < 1:
        denominator = round(1 / numeric)
        return f"1/{denominator}"
    return f"{numeric:.1f}".rstrip("0").rstrip(".")


def get_image_size(path: Path) -> tuple[int, int]:
    with Image.open(path) as image:
        image = ImageOps.exif_transpose(image)
        return image.size


def parse_color(value: str, fallback: tuple[int, int, int]) -> tuple[int, int, int]:
    if not isinstance(value, str):
        return fallback
    match = re.fullmatch(r"#?([0-9a-fA-F]{6})", value.strip())
    if not match:
        return fallback
    hex_value = match.group(1)
    return tuple(int(hex_value[index : index + 2], 16) for index in (0, 2, 4))


def merged_settings(settings: dict[str, Any] | None) -> dict[str, Any]:
    merged = dict(DEFAULT_SETTINGS)
    if settings:
        merged.update(settings)
    if not merged.get("_fixedSlideSize"):
        apply_derived_layout_settings(merged)
    return merged


def preview_cache_key(album_id: str, index: int, photo: dict[str, Any], settings: dict[str, Any], max_size: int) -> str:
    path = Path(clean_text(photo.get("path")))
    try:
        stat = path.stat()
        file_signature = [str(path), stat.st_mtime_ns, stat.st_size]
    except OSError:
        file_signature = [str(path), 0, 0]
    signature = {
        "albumId": album_id,
        "index": index,
        "file": file_signature,
        "maxSize": max_size,
        "logo": render_logo_signature(photo.get("exif") or {}, settings),
        "settings": settings,
    }
    payload = json.dumps(signature, sort_keys=True, ensure_ascii=False, default=str).encode("utf-8")
    return hashlib.sha256(payload).hexdigest()


def get_cached_preview(key: str) -> bytes | None:
    with PREVIEW_CACHE_LOCK:
        data = PREVIEW_CACHE.get(key)
        if data is None:
            return None
        PREVIEW_CACHE.move_to_end(key)
        return data


def remember_preview(key: str, data: bytes) -> None:
    global PREVIEW_CACHE_BYTES
    with PREVIEW_CACHE_LOCK:
        old_data = PREVIEW_CACHE.pop(key, None)
        if old_data is not None:
            PREVIEW_CACHE_BYTES -= len(old_data)
        PREVIEW_CACHE[key] = data
        PREVIEW_CACHE_BYTES += len(data)
        PREVIEW_CACHE.move_to_end(key)
        byte_limit = preview_cache_byte_limit()
        item_limit = preview_cache_item_limit(byte_limit)
        while len(PREVIEW_CACHE) > 1 and (
            len(PREVIEW_CACHE) > item_limit or PREVIEW_CACHE_BYTES > byte_limit
        ):
            _, removed = PREVIEW_CACHE.popitem(last=False)
            PREVIEW_CACHE_BYTES -= len(removed)


def preview_cache_byte_limit() -> int:
    available = available_memory_bytes()
    if available is None:
        return PREVIEW_CACHE_DEFAULT_BYTES
    if available < 2 * GIB:
        return int(max(PREVIEW_CACHE_MIN_BYTES, min(512 * MIB, available // 8)))
    budget = (available - GIB) // 3
    return int(max(512 * MIB, min(PREVIEW_CACHE_MAX_BYTES, budget)))


def preview_cache_item_limit(byte_limit: int) -> int:
    estimated = max(1, byte_limit // PREVIEW_CACHE_ESTIMATED_ITEM_BYTES)
    return int(max(PREVIEW_CACHE_MIN_ITEMS, min(PREVIEW_CACHE_MAX_ITEMS, estimated)))


def available_memory_bytes() -> int | None:
    if os.name == "nt":
        return windows_available_memory_bytes()
    try:
        page_size = os.sysconf("SC_PAGE_SIZE")
        available_pages = os.sysconf("SC_AVPHYS_PAGES")
        return int(page_size * available_pages)
    except (AttributeError, OSError, ValueError):
        return None


def windows_available_memory_bytes() -> int | None:
    try:
        import ctypes

        class MemoryStatusEx(ctypes.Structure):
            _fields_ = [
                ("dwLength", ctypes.c_ulong),
                ("dwMemoryLoad", ctypes.c_ulong),
                ("ullTotalPhys", ctypes.c_ulonglong),
                ("ullAvailPhys", ctypes.c_ulonglong),
                ("ullTotalPageFile", ctypes.c_ulonglong),
                ("ullAvailPageFile", ctypes.c_ulonglong),
                ("ullTotalVirtual", ctypes.c_ulonglong),
                ("ullAvailVirtual", ctypes.c_ulonglong),
                ("ullAvailExtendedVirtual", ctypes.c_ulonglong),
            ]

        status = MemoryStatusEx()
        status.dwLength = ctypes.sizeof(MemoryStatusEx)
        if ctypes.windll.kernel32.GlobalMemoryStatusEx(ctypes.byref(status)):
            return int(status.ullAvailPhys)
    except Exception:
        return None
    return None


def apply_derived_layout_settings(settings: dict[str, Any]) -> None:
    ratio_width, ratio_height = parse_aspect_ratio(clean_text(settings.get("slideAspectRatio")), (16, 9))
    long_edge = clamp_int(settings.get("slideLongEdge") or settings.get("slideWidth"), 800, 8000, 1920)
    if ratio_width >= ratio_height:
        settings["slideWidth"] = long_edge
        settings["slideHeight"] = max(1, round(long_edge * ratio_height / ratio_width))
    else:
        settings["slideHeight"] = long_edge
        settings["slideWidth"] = max(1, round(long_edge * ratio_width / ratio_height))

    if "pageMarginXPct" not in settings and settings.get("safeWidthPct") not in (None, ""):
        settings["pageMarginXPct"] = max(0, (100 - float(settings["safeWidthPct"])) / 2)
    if "pageMarginYPct" not in settings and settings.get("safeHeightPct") not in (None, ""):
        banner = float(settings.get("bannerHeightPct", DEFAULT_SETTINGS["bannerHeightPct"]))
        gap = float(settings.get("gapPct", DEFAULT_SETTINGS["gapPct"]))
        settings["pageMarginYPct"] = max(0, (100 - float(settings["safeHeightPct"]) - banner - gap) / 2)

    settings["safeWidthPct"] = max(1, 100 - 2 * float(settings.get("pageMarginXPct", 2.5)))
    settings["safeHeightPct"] = max(1, 100 - 2 * float(settings.get("pageMarginYPct", 1.85)))


def parse_aspect_ratio(value: str, fallback: tuple[int, int]) -> tuple[int, int]:
    match = re.fullmatch(r"\s*(\d+(?:\.\d+)?)\s*[:/]\s*(\d+(?:\.\d+)?)\s*", value or "")
    if not match:
        return fallback
    width = float(match.group(1))
    height = float(match.group(2))
    if width <= 0 or height <= 0:
        return fallback
    return max(1, round(width * 1000)), max(1, round(height * 1000))


def render_composite(photo: dict[str, Any], settings: dict[str, Any]) -> Image.Image:
    settings = merged_settings(settings)
    image = load_oriented_rgb(photo["path"])

    layout = composite_layout(image.width, image.height, settings)
    return render_composite_from_image(photo, settings, image, layout)


def render_preview_composite(photo: dict[str, Any], settings: dict[str, Any], max_size: int) -> Image.Image:
    settings = preview_render_settings(settings, max_size)
    image_width, image_height = oriented_image_size(photo)
    layout = composite_layout(image_width, image_height, settings)
    image = load_oriented_rgb(photo["path"], (layout["imageWidth"], layout["imageHeight"]))
    return render_composite_from_image(photo, settings, image, layout)


def preview_render_settings(settings: dict[str, Any], max_size: int) -> dict[str, Any]:
    merged = merged_settings(settings)
    max_size = clamp_int(max_size, 600, 2400, 1600)
    long_edge = max(int(merged["slideWidth"]), int(merged["slideHeight"]))
    if long_edge <= max_size:
        return merged
    scale = max_size / long_edge
    preview = dict(merged)
    preview["slideWidth"] = max(1, round(int(merged["slideWidth"]) * scale))
    preview["slideHeight"] = max(1, round(int(merged["slideHeight"]) * scale))
    preview["slideLongEdge"] = max(preview["slideWidth"], preview["slideHeight"])
    preview["_fixedSlideSize"] = True
    return preview


def render_composite_from_image(
    photo: dict[str, Any],
    settings: dict[str, Any],
    image: Image.Image,
    layout: dict[str, int],
) -> Image.Image:
    slide_width = layout["slideWidth"]
    slide_height = layout["slideHeight"]
    background = parse_color(settings.get("background", ""), (246, 244, 239))
    canvas = Image.new("RGB", (slide_width, slide_height), background)
    image = image.resize((layout["imageWidth"], layout["imageHeight"]), RESAMPLE)

    if settings.get("shadow", True):
        add_shadow(canvas, image, layout["imageLeft"], layout["imageTop"])
    canvas.paste(image, (layout["imageLeft"], layout["imageTop"]))

    draw_banner(
        canvas,
        photo,
        settings,
        layout["bannerLeft"],
        layout["bannerTop"],
        layout["bannerWidth"],
        layout["bannerHeight"],
    )
    return canvas


def load_oriented_rgb(path: str | Path, target_size: tuple[int, int] | None = None) -> Image.Image:
    with Image.open(path) as source_image:
        apply_decode_draft(source_image, target_size)
        image = ImageOps.exif_transpose(source_image)
        if image.mode != "RGB":
            return image.convert("RGB")
        return image.copy()


def apply_decode_draft(image: Image.Image, target_size: tuple[int, int] | None) -> None:
    if not target_size:
        return
    try:
        target_width = max(1, int(target_size[0]))
        target_height = max(1, int(target_size[1]))
        orientation = image.getexif().get(274)
        if orientation in {5, 6, 7, 8}:
            target_width, target_height = target_height, target_width
        image.draft("RGB", (target_width * 2, target_height * 2))
    except Exception:
        pass


def composite_layout(image_width: int, image_height: int, settings: dict[str, Any]) -> dict[str, int]:
    fixed_size = bool(settings.get("_fixedSlideSize"))
    max_edge = EXPORT_MAX_EDGE if fixed_size else 8000
    slide_width = clamp_int(settings.get("slideWidth"), 1 if fixed_size else 800, max_edge, 1920)
    slide_height = clamp_int(settings.get("slideHeight"), 1 if fixed_size else 600, max_edge, 1080)
    banner_height = max(24, int(slide_height * float(settings["bannerHeightPct"]) / 100))
    gap = max(0, int(slide_height * float(settings["gapPct"]) / 100))
    margin_x = max(0, int(slide_width * float(settings.get("pageMarginXPct", 2.5)) / 100))
    margin_y = max(0, int(slide_height * float(settings.get("pageMarginYPct", 1.85)) / 100))
    max_width = max(1, slide_width - margin_x * 2)
    max_group_height = max(1, slide_height - margin_y * 2)
    max_image_height = max(1, max_group_height - banner_height - gap)
    scale = min(max_width / image_width, max_image_height / image_height)
    image_scale_cap = settings.get("_imageScaleCap")
    if image_scale_cap:
        scale = min(scale, max(0.01, float(image_scale_cap)))
    fitted_image_width = max(1, round(image_width * scale))
    fitted_image_height = max(1, round(image_height * scale))

    total_height = fitted_image_height + gap + banner_height
    image_left = (slide_width - fitted_image_width) // 2
    image_top = (slide_height - total_height) // 2

    banner_width = max(120, int(slide_width * float(settings["bannerWidthPct"]) / 100))
    banner_left = (slide_width - banner_width) // 2
    banner_top = image_top + fitted_image_height + gap
    return {
        "slideWidth": slide_width,
        "slideHeight": slide_height,
        "imageLeft": image_left,
        "imageTop": image_top,
        "imageWidth": fitted_image_width,
        "imageHeight": fitted_image_height,
        "bannerLeft": banner_left,
        "bannerTop": banner_top,
        "bannerWidth": banner_width,
        "bannerHeight": banner_height,
    }


def banner_edit_layout(photo: dict[str, Any], settings: dict[str, Any]) -> dict[str, Any]:
    settings = merged_settings(settings)
    image_width, image_height = oriented_image_size(photo)
    layout = composite_layout(image_width, image_height, settings)
    fields = banner_edit_fields(photo, settings, layout)
    return {
        "slideWidth": layout["slideWidth"],
        "slideHeight": layout["slideHeight"],
        "banner": rect_payload(
            layout["bannerLeft"],
            layout["bannerTop"],
            layout["bannerWidth"],
            layout["bannerHeight"],
            layout["slideWidth"],
            layout["slideHeight"],
        ),
        "fields": fields,
    }


def export_settings_for_photo(photo: dict[str, Any], settings: dict[str, Any]) -> dict[str, Any]:
    merged = merged_settings(settings)
    scale_pct = clamp_int(merged.get("exportScalePct"), 10, 100, 100)
    scale = scale_pct / 100
    image_width, image_height = oriented_image_size(photo)

    ratio_width, ratio_height = parse_aspect_ratio(clean_text(merged.get("slideAspectRatio")), (16, 9))
    aspect = ratio_width / ratio_height
    width_factor = max(0.01, 1 - 2 * float(merged.get("pageMarginXPct", 2.5)) / 100)
    height_factor = max(
        0.01,
        1
        - 2 * float(merged.get("pageMarginYPct", 1.85)) / 100
        - float(merged.get("bannerHeightPct", DEFAULT_SETTINGS["bannerHeightPct"])) / 100
        - float(merged.get("gapPct", DEFAULT_SETTINGS["gapPct"])) / 100,
    )
    required_height = max(
        image_width * scale / (aspect * width_factor),
        image_height * scale / height_factor,
    )
    slide_height = max(1, math.ceil(required_height))
    slide_width = max(1, math.ceil(slide_height * aspect))

    fixed = dict(merged)
    fixed.update(
        {
            "slideWidth": slide_width,
            "slideHeight": slide_height,
            "slideLongEdge": max(slide_width, slide_height),
            "exportScalePct": scale_pct,
            "_fixedSlideSize": True,
            "_imageScaleCap": scale,
        }
    )
    validate_export_canvas(fixed)
    return fixed


def export_settings_for_deck(photos: list[dict[str, Any]], settings: dict[str, Any]) -> dict[str, Any]:
    if not photos:
        return export_settings_for_photo({"path": "", "width": 1, "height": 1}, settings)
    per_photo = [export_settings_for_photo(photo, settings) for photo in photos]
    slide_height = max(int(item["slideHeight"]) for item in per_photo)
    first = per_photo[0]
    aspect = int(first["slideWidth"]) / max(1, int(first["slideHeight"]))
    fixed = dict(first)
    fixed["slideHeight"] = slide_height
    fixed["slideWidth"] = max(1, math.ceil(slide_height * aspect))
    fixed["slideLongEdge"] = max(int(fixed["slideWidth"]), int(fixed["slideHeight"]))
    validate_export_canvas(fixed)
    return fixed


def oriented_image_size(photo: dict[str, Any]) -> tuple[int, int]:
    path_text = clean_text(photo.get("path"))
    if path_text:
        try:
            with Image.open(path_text) as image:
                width, height = image.size
                orientation = image.getexif().get(274)
                if orientation in {5, 6, 7, 8}:
                    return height, width
                return width, height
        except Exception:
            pass
    return (
        max(1, clamp_int(photo.get("width"), 1, EXPORT_MAX_EDGE, 1)),
        max(1, clamp_int(photo.get("height"), 1, EXPORT_MAX_EDGE, 1)),
    )


def validate_export_canvas(settings: dict[str, Any]) -> None:
    width = int(settings["slideWidth"])
    height = int(settings["slideHeight"])
    if width > EXPORT_MAX_EDGE or height > EXPORT_MAX_EDGE:
        raise ValueError(f"导出尺寸 {width}x{height} 超过上限，请降低导出分辨率百分比。")
    if width * height > EXPORT_MAX_PIXELS:
        raise ValueError(f"导出尺寸 {width}x{height} 过大，请降低导出分辨率百分比。")


def clamp_int(value: Any, minimum: int, maximum: int, fallback: int) -> int:
    try:
        return max(minimum, min(maximum, int(value)))
    except Exception:
        return fallback


def add_shadow(canvas: Image.Image, image: Image.Image, left: int, top: int) -> None:
    shadow = Image.new("RGBA", canvas.size, (0, 0, 0, 0))
    shadow_draw = ImageDraw.Draw(shadow)
    shadow_draw.rectangle(
        [left + 8, top + 12, left + image.width + 8, top + image.height + 12],
        fill=(0, 0, 0, 90),
    )
    shadow = shadow.filter(ImageFilter.GaussianBlur(18))
    canvas.paste(Image.alpha_composite(canvas.convert("RGBA"), shadow).convert("RGB"))


def draw_banner(
    canvas: Image.Image,
    photo: dict[str, Any],
    settings: dict[str, Any],
    left: int,
    top: int,
    width: int,
    height: int,
) -> None:
    exif = photo.get("exif", {})
    if not isinstance(exif, dict):
        exif = {}
    text_overrides = banner_text_overrides(photo, settings)
    overlay = Image.new("RGBA", canvas.size, (0, 0, 0, 0))
    overlay_draw = ImageDraw.Draw(overlay)
    color = parse_color(settings.get("bannerColor", ""), (0, 0, 0))
    opacity = int(max(0, min(1, float(settings.get("bannerOpacity", 0.6)))) * 255)
    overlay_draw.rectangle([left, top, left + width, top + height], fill=(*color, opacity))
    canvas.paste(Image.alpha_composite(canvas.convert("RGBA"), overlay).convert("RGB"))

    draw = ImageDraw.Draw(canvas)
    macro_unit = width / 600
    horizontal_margin = max(1, round(15 * macro_unit))
    logo_gap = max(1, round(10 * macro_unit))
    text_margin = max(0, round(7.2 * macro_unit))
    logo_left = left + horizontal_margin
    logo_height = max(1, round(height * 0.75))
    logo_top = top + (height - logo_height) // 2
    logo_box_width = max(1, round(width * 0.22))
    content_right = left + width - horizontal_margin - text_margin

    logo_width = draw_logo_or_brand(
        canvas,
        draw,
        exif,
        settings,
        logo_left,
        logo_top,
        logo_box_width,
        logo_height,
        clean_text(text_overrides.get("brand")),
    )
    content_left = logo_left + logo_width + logo_gap + text_margin

    model = clean_text(text_overrides.get("model")) or exif.get("model") or "Unknown camera"
    lens = clean_text(text_overrides.get("lens")) or exif.get("lens") or "Unknown lens"
    info_font_size = max(10, int(canvas.height * float(settings["infoFontPct"]) / 100))
    info_font = load_font(info_font_size, bold=True, text=f"{model}\n{lens}")
    info_width = max(80, round(width * 0.5) - text_margin * 2)
    info_lines = [
        ellipsize(draw, model, info_font, info_width),
        ellipsize(draw, lens, info_font, info_width),
    ]
    info_line_height = font_line_height(info_font)
    info_y = top + (height - info_line_height * len(info_lines)) / 2
    for line_index, line in enumerate(info_lines):
        draw.text(
            (content_left, info_y + line_index * info_line_height),
            line,
            font=info_font,
            fill=(224, 224, 224),
        )

    params = clean_text(text_overrides.get("params")) or format_params(exif)
    param_font_size = max(12, int(canvas.height * float(settings["paramFontPct"]) / 100))
    param_font = load_font(param_font_size, bold=True, text=params)
    param_left = left + int(width * 0.5)
    param_width = max(80, content_right - param_left)
    param_font = fit_font(params, param_font, param_width, int(height * 0.72))
    params = ellipsize(draw, params, param_font, param_width)
    bbox = draw.textbbox((0, 0), params, font=param_font)
    param_x = content_right - (bbox[2] - bbox[0])
    param_y = top + (height - font_line_height(param_font)) / 2
    draw.text((param_x, param_y), params, font=param_font, fill=(255, 255, 255))


def banner_edit_fields(
    photo: dict[str, Any],
    settings: dict[str, Any],
    layout: dict[str, int],
) -> dict[str, dict[str, Any]]:
    exif = photo.get("exif", {})
    if not isinstance(exif, dict):
        exif = {}
    text_overrides = banner_text_overrides(photo, settings)
    left = layout["bannerLeft"]
    top = layout["bannerTop"]
    width = layout["bannerWidth"]
    height = layout["bannerHeight"]
    slide_width = layout["slideWidth"]
    slide_height = layout["slideHeight"]

    measure_draw = ImageDraw.Draw(Image.new("RGB", (1, 1)))
    macro_unit = width / 600
    horizontal_margin = max(1, round(15 * macro_unit))
    logo_gap = max(1, round(10 * macro_unit))
    text_margin = max(0, round(7.2 * macro_unit))
    logo_left = left + horizontal_margin
    logo_height = max(1, round(height * 0.75))
    logo_top = top + (height - logo_height) // 2
    logo_box_width = max(1, round(width * 0.22))
    content_right = left + width - horizontal_margin - text_margin

    logo_width, has_logo, brand_font_size, brand_line_height = logo_or_brand_metrics(
        measure_draw,
        exif,
        settings,
        layout,
        logo_box_width,
        logo_height,
        clean_text(text_overrides.get("brand")),
    )
    content_left = logo_left + logo_width + logo_gap + text_margin

    model = clean_text(text_overrides.get("model")) or exif.get("model") or "Unknown camera"
    lens = clean_text(text_overrides.get("lens")) or exif.get("lens") or "Unknown lens"
    info_font_size = max(10, int(slide_height * float(settings["infoFontPct"]) / 100))
    info_font = load_font(info_font_size, bold=True, text=f"{model}\n{lens}")
    info_line_height = font_line_height(info_font)
    info_y = top + (height - info_line_height * 2) / 2

    info_width = max(80, round(width * 0.5) - text_margin * 2)
    param_left = left + int(width * 0.5)
    param_width = max(80, content_right - param_left)
    visual_info_width = max(1, min(info_width, param_left - content_left - text_margin))

    params = clean_text(text_overrides.get("params")) or format_params(exif)
    param_font_size = max(12, int(slide_height * float(settings["paramFontPct"]) / 100))
    param_font = load_font(param_font_size, bold=True, text=params)
    param_font = fit_font(params, param_font, param_width, int(height * 0.72))
    param_line_height = font_line_height(param_font)
    param_y = top + (height - param_line_height) / 2

    fields: dict[str, dict[str, Any]] = {
        "brand": rect_payload(
            logo_left,
            logo_top,
            max(1, logo_width),
            logo_height,
            slide_width,
            slide_height,
            hidden=has_logo,
            font_size=brand_font_size,
            line_height=brand_line_height,
            align="left",
        ),
        "model": rect_payload(
            content_left,
            info_y,
            visual_info_width,
            info_line_height,
            slide_width,
            slide_height,
            font_size=info_font_size,
            line_height=info_line_height,
            align="left",
        ),
        "lens": rect_payload(
            content_left,
            info_y + info_line_height,
            visual_info_width,
            info_line_height,
            slide_width,
            slide_height,
            font_size=info_font_size,
            line_height=info_line_height,
            align="left",
        ),
        "params": rect_payload(
            param_left,
            param_y,
            param_width,
            param_line_height,
            slide_width,
            slide_height,
            font_size=getattr(param_font, "size", param_font_size),
            line_height=param_line_height,
            align="right",
        ),
    }
    return fields


def logo_or_brand_metrics(
    draw: ImageDraw.ImageDraw,
    exif: dict[str, str],
    settings: dict[str, Any],
    layout: dict[str, int],
    width: int,
    height: int,
    brand_override: str = "",
) -> tuple[int, bool, int, int]:
    logo_path = render_logo_path(exif, settings)
    logo = load_logo_image(str(logo_path) if logo_path else "")
    if logo:
        try:
            scale = min(width / logo.width, height / logo.height)
            return max(1, int(logo.width * scale)), True, 0, height
        except Exception:
            pass

    brand = brand_override or clean_text(settings.get("brandText")) or infer_brand(exif)
    font_size = max(12, int(layout["slideHeight"] * float(settings["brandFontPct"]) / 100))
    font = fit_font(brand, load_font(font_size, bold=True, text=brand), width, height)
    brand = ellipsize(draw, brand, font, width)
    bbox = draw.textbbox((0, 0), brand, font=font)
    return max(1, bbox[2] - bbox[0]), False, getattr(font, "size", font_size), font_line_height(font)


def rect_payload(
    left: float,
    top: float,
    width: float,
    height: float,
    slide_width: int,
    slide_height: int,
    *,
    hidden: bool = False,
    font_size: int | float | None = None,
    line_height: int | float | None = None,
    align: str = "left",
) -> dict[str, Any]:
    payload: dict[str, Any] = {
        "left": left,
        "top": top,
        "width": width,
        "height": height,
        "leftPct": left / slide_width * 100,
        "topPct": top / slide_height * 100,
        "widthPct": width / slide_width * 100,
        "heightPct": height / slide_height * 100,
        "hidden": hidden,
        "align": align,
    }
    if font_size is not None:
        payload["fontSize"] = font_size
    if line_height is not None:
        payload["lineHeight"] = line_height
    return payload


def banner_text_overrides(photo: dict[str, Any], settings: dict[str, Any]) -> dict[str, str]:
    overrides = settings.get("bannerTextOverrides")
    if not isinstance(overrides, dict):
        return {}
    name = clean_text(photo.get("name"))
    size = clean_text(photo.get("size"))
    keys = [
        photo_override_key(photo),
        f"{name}:{size}" if name or size else "",
        name,
    ]
    path_text = clean_text(photo.get("path"))
    if path_text:
        file_name = Path(path_text).name
        keys.extend([
            f"{file_name}:{size}" if file_name or size else "",
            file_name,
        ])
    for key in keys:
        if not key:
            continue
        values = overrides.get(key)
        if isinstance(values, dict):
            return {
                field: clean_text(values.get(field))
                for field in ("brand", "model", "lens", "params")
                if clean_text(values.get(field))
            }
    return {}


def photo_override_key(photo: dict[str, Any]) -> str:
    path_text = clean_text(photo.get("path"))
    if path_text:
        return path_text
    return f"{clean_text(photo.get('name'))}:{clean_text(photo.get('size'))}"


def draw_logo_or_brand(
    canvas: Image.Image,
    draw: ImageDraw.ImageDraw,
    exif: dict[str, str],
    settings: dict[str, Any],
    left: int,
    top: int,
    width: int,
    height: int,
    brand_override: str = "",
) -> int:
    logo_path = render_logo_path(exif, settings)
    logo = load_logo_image(str(logo_path) if logo_path else "")
    if logo:
        try:
            scale = min(width / logo.width, height / logo.height)
            logo_size = (max(1, int(logo.width * scale)), max(1, int(logo.height * scale)))
            logo = logo.resize(logo_size, RESAMPLE)
            canvas.paste(logo, (left, top + (height - logo.height) // 2), logo)
            return logo.width
        except Exception:
            pass

    brand = brand_override or clean_text(settings.get("brandText")) or infer_brand(exif)
    font_size = max(12, int(canvas.height * float(settings["brandFontPct"]) / 100))
    font = fit_font(brand, load_font(font_size, bold=True, text=brand), width, height)
    brand = ellipsize(draw, brand, font, width)
    bbox = draw.textbbox((0, 0), brand, font=font)
    draw.text(
        (left, top + (height - (bbox[3] - bbox[1])) / 2 - bbox[1]),
        brand,
        font=font,
        fill=(255, 255, 255),
    )
    return bbox[2] - bbox[0]


def render_logo_path(exif: dict[str, str], settings: dict[str, Any]) -> Path | None:
    manual_path = clean_text(settings.get("logoPath"))
    if manual_path:
        path = Path(manual_path).expanduser()
        if path.exists() and path.is_file() and path.suffix.lower() in LOGO_EXTENSIONS:
            return path
    return matched_builtin_logo_path(exif)


def render_logo_signature(exif: dict[str, str], settings: dict[str, Any]) -> list[Any]:
    path = render_logo_path(exif, settings)
    if not path:
        return ["none", logo_rules_mtime()]
    try:
        stat = path.stat()
        return [str(path.resolve()), stat.st_mtime_ns, stat.st_size, logo_rules_mtime()]
    except OSError:
        return [str(path), 0, 0, logo_rules_mtime()]


def load_logo_image(logo_path: str) -> Image.Image | None:
    path_text = clean_text(logo_path)
    if not path_text:
        return None
    path = Path(path_text)
    if not path.exists() or not path.is_file() or path.suffix.lower() not in LOGO_EXTENSIONS:
        return None
    try:
        stat = path.stat()
        cache_key = str(path.resolve())
        signature = (stat.st_mtime_ns, stat.st_size)
        with LOGO_CACHE_LOCK:
            cached = LOGO_CACHE.get(cache_key)
            if cached and cached[0] == signature:
                return cached[1].copy()

        with Image.open(path) as image:
            logo = ImageOps.exif_transpose(image).convert("RGBA")

        with LOGO_CACHE_LOCK:
            LOGO_CACHE[cache_key] = (signature, logo.copy())
            while len(LOGO_CACHE) > 8:
                LOGO_CACHE.pop(next(iter(LOGO_CACHE)))
        return logo
    except Exception:
        return None


def matched_builtin_logo_path(exif: dict[str, str]) -> Path | None:
    text = normalized_logo_text(" ".join([
        clean_text(exif.get("make")),
        clean_text(exif.get("model")),
    ]))
    if not text:
        return None
    for rule in load_logo_rules():
        terms = rule.get("match") or []
        if any(normalized_logo_text(term) and normalized_logo_text(term) in text for term in terms):
            file_name = clean_text(rule.get("file"))
            path = (LOGO_DIR / file_name).resolve()
            try:
                if path.is_file() and str(path).startswith(str(LOGO_DIR.resolve())):
                    return path
            except OSError:
                return None
    return None


def load_logo_rules() -> list[dict[str, Any]]:
    global LOGO_RULES_CACHE
    try:
        mtime = LOGO_RULES_FILE.stat().st_mtime_ns
    except OSError:
        return []
    with LOGO_RULES_LOCK:
        if LOGO_RULES_CACHE and LOGO_RULES_CACHE[0] == mtime:
            return LOGO_RULES_CACHE[1]
        try:
            data = json.loads(LOGO_RULES_FILE.read_text(encoding="utf-8"))
            rules = data.get("logos", [])
            if not isinstance(rules, list):
                rules = []
        except Exception:
            rules = []
        LOGO_RULES_CACHE = (mtime, rules)
        return rules


def logo_rules_mtime() -> int:
    try:
        return LOGO_RULES_FILE.stat().st_mtime_ns
    except OSError:
        return 0


def normalized_logo_text(value: Any) -> str:
    return re.sub(r"[^a-z0-9]+", " ", clean_text(value).lower()).strip()


def infer_brand(exif: dict[str, str]) -> str:
    text = (exif.get("make") or exif.get("model") or "CAMERA").strip()
    if not text:
        return "CAMERA"
    return text.split()[0].upper()


def format_params(exif: dict[str, str]) -> str:
    values: list[str] = []
    if exif.get("focalLength"):
        values.append(exif["focalLength"])
    if exif.get("fNumber"):
        values.append(f"F{exif['fNumber']}")
    if exif.get("exposureTime"):
        values.append(f"{exif['exposureTime']}s")
    if exif.get("iso"):
        values.append(f"ISO {exif['iso']}")
    return "    ".join(values) or "No EXIF"


def load_font(size: int, bold: bool = False, text: str = "") -> ImageFont.FreeTypeFont | ImageFont.ImageFont:
    for candidate in font_candidates(bold, text):
        if candidate.exists():
            try:
                return ImageFont.truetype(str(candidate), size)
            except OSError:
                continue
    return ImageFont.load_default()


def font_candidates(bold: bool, text: str = "") -> list[Path]:
    font_dir = Path(os.environ.get("WINDIR", "C:/Windows")) / "Fonts"
    preferred = [
        "segoeuib.ttf" if bold else "segoeui.ttf",
        "arialbd.ttf" if bold else "arial.ttf",
    ]
    multilingual = [
        "NotoSansSC-VF.ttf",
        "msyhbd.ttc" if bold else "msyh.ttc",
        "msyh.ttc",
        "msyhbd.ttc",
        "simhei.ttf",
        "simsun.ttc",
        "msjhbd.ttc" if bold else "msjh.ttc",
        "msjh.ttc",
        "YuGothB.ttc" if bold else "YuGothR.ttc",
        "YuGothM.ttc",
        "malgunbd.ttf" if bold else "malgun.ttf",
        "malgun.ttf",
        "seguisym.ttf",
    ]
    if needs_multilingual_font(text):
        names = multilingual + preferred
    else:
        names = preferred + multilingual

    candidates: list[Path] = []
    seen: set[str] = set()
    for name in names:
        path = font_dir / name
        key = str(path).lower()
        if key not in seen:
            candidates.append(path)
            seen.add(key)
    return candidates


def needs_multilingual_font(text: str) -> bool:
    return any(
        "\u2e80" <= char <= "\u9fff"
        or "\uf900" <= char <= "\ufaff"
        or "\u3040" <= char <= "\u30ff"
        or "\uac00" <= char <= "\ud7af"
        for char in text
    )


def font_line_height(font: ImageFont.ImageFont) -> int:
    if isinstance(font, ImageFont.FreeTypeFont):
        ascent, descent = font.getmetrics()
        return ascent + descent
    bbox = ImageDraw.Draw(Image.new("RGB", (1, 1))).textbbox((0, 0), "Ag", font=font)
    return bbox[3] - bbox[1]


def fit_font(text: str, font: ImageFont.ImageFont, max_width: int, max_height: int) -> ImageFont.ImageFont:
    if not isinstance(font, ImageFont.FreeTypeFont):
        return font
    size = font.size
    while size > 8:
        test_font = load_font(size, bold=True, text=text)
        bbox = ImageDraw.Draw(Image.new("RGB", (1, 1))).textbbox((0, 0), text, font=test_font)
        if bbox[2] - bbox[0] <= max_width and bbox[3] - bbox[1] <= max_height:
            return test_font
        size -= 1
    return load_font(8, bold=True, text=text)


def ellipsize(draw: ImageDraw.ImageDraw, text: str, font: ImageFont.ImageFont, max_width: int) -> str:
    if draw.textbbox((0, 0), text, font=font)[2] <= max_width:
        return text
    suffix = "..."
    while text:
        candidate = text[:-1].rstrip() + suffix
        if draw.textbbox((0, 0), candidate, font=font)[2] <= max_width:
            return candidate
        text = text[:-1]
    return suffix


def scan_album(folder: Path, recursive: bool, sort_mode: str = "dateAsc") -> dict[str, Any]:
    if not folder.exists() or not folder.is_dir():
        raise ValueError("Folder does not exist.")
    sort_mode = normalize_sort_mode(sort_mode)
    paths = list_images(folder, recursive)
    exif_map, exif_source = read_exif_batch(paths, folder, recursive)
    photos: list[dict[str, Any]] = []
    for index, path in enumerate(paths):
        metadata = dict(exif_map.get(str(path.resolve()), exif_map.get(str(path), {})))
        width = clamp_int(metadata.pop("imageWidth", 0), 0, 100000, 0)
        height = clamp_int(metadata.pop("imageHeight", 0), 0, 100000, 0)
        if not width or not height:
            try:
                width, height = get_image_size(path)
            except Exception:
                width, height = 0, 0
        photo = Photo(
            index=index,
            path=str(path),
            name=path.name,
            size=path.stat().st_size,
            width=width,
            height=height,
            exif=metadata,
        )
        photos.append(photo.__dict__)
    photos = sort_photos(photos, sort_mode)
    for index, photo in enumerate(photos):
        photo["index"] = index
    album_id = uuid.uuid4().hex
    SESSIONS[album_id] = {
        "root": str(folder),
        "photos": photos,
        "created": time.time(),
        "recursive": recursive,
        "sortMode": sort_mode,
        "exifSource": exif_source,
    }
    return {
        "albumId": album_id,
        "root": str(folder),
        "count": len(photos),
        "recursive": recursive,
        "sortMode": sort_mode,
        "exifSource": exif_source,
        "photos": public_photos(photos),
        "settings": DEFAULT_SETTINGS,
    }


def sort_photos(photos: list[dict[str, Any]], sort_mode: str) -> list[dict[str, Any]]:
    if sort_mode == "nameAsc":
        return sorted(photos, key=lambda photo: natural_key(str(photo.get("name", ""))))
    if sort_mode == "nameDesc":
        return sorted(photos, key=lambda photo: natural_key(str(photo.get("name", ""))), reverse=True)
    descending = sort_mode == "dateDesc"
    return sorted(photos, key=lambda photo: photo_time_key(photo, descending))


def photo_time_key(photo: dict[str, Any], descending: bool = False) -> tuple[int, float, list[Any]]:
    timestamp = exif_timestamp(photo.get("exif", {}).get("dateTime"))
    if timestamp is None:
        try:
            timestamp = Path(str(photo.get("path", ""))).stat().st_mtime
        except OSError:
            timestamp = None
    if timestamp is None:
        return (1, 0.0, natural_key(str(photo.get("name", ""))))
    return (0, -timestamp if descending else timestamp, natural_key(str(photo.get("name", ""))))


def exif_timestamp(value: Any) -> float | None:
    text = clean_text(value).strip()
    if not text:
        return None
    text = re.sub(r"\s+[A-Z]{2,5}$", "", text)
    candidates = [text, text.replace(":", "-", 2)]
    for candidate in candidates:
        candidate = candidate.replace("T", " ")
        candidate = re.sub(r"([+-]\d{2}):?(\d{2})$", r"\1:\2", candidate)
        try:
            return datetime.fromisoformat(candidate).timestamp()
        except ValueError:
            pass
        for fmt in ("%Y:%m:%d %H:%M:%S", "%Y-%m-%d %H:%M:%S", "%Y:%m:%d %H:%M", "%Y-%m-%d %H:%M"):
            try:
                return datetime.strptime(candidate, fmt).timestamp()
            except ValueError:
                pass
    return None


def public_photos(photos: list[dict[str, Any]]) -> list[dict[str, Any]]:
    return [
        {
            "index": photo["index"],
            "name": photo["name"],
            "size": photo["size"],
            "width": photo["width"],
            "height": photo["height"],
            "exif": photo["exif"],
        }
        for photo in photos
    ]


def selected_photos(album: dict[str, Any], selection: list[int] | None) -> list[dict[str, Any]]:
    photos = album["photos"]
    if not selection:
        return photos
    indexes = set(int(value) for value in selection)
    return [photo for photo in photos if int(photo["index"]) in indexes]


def export_images(album_id: str, payload: dict[str, Any]) -> dict[str, Any]:
    album = get_album(album_id)
    settings = merged_settings(payload.get("settings"))
    photos = selected_photos(album, payload.get("selection"))
    root = Path(album["root"])
    output_dir = Path(payload.get("outputDir") or root / "banner_output" / f"images_{timestamp()}")
    output_dir.mkdir(parents=True, exist_ok=True)

    image_ext = export_extension(settings)
    targets = unique_export_paths(photos, output_dir, image_ext, root)
    written = export_photo_images_parallel(photos, targets, settings)
    return {"count": len(written), "outputDir": str(output_dir), "files": written[:12]}


def export_pptx(album_id: str, payload: dict[str, Any]) -> dict[str, Any]:
    album = get_album(album_id)
    settings = merged_settings(payload.get("settings"))
    photos = selected_photos(album, payload.get("selection"))
    root = Path(album["root"])
    output_dir = Path(payload.get("outputDir") or root / "banner_output")
    output_dir.mkdir(parents=True, exist_ok=True)
    output_file = Path(payload.get("outputFile") or output_dir / f"{root.name}_banner_{timestamp()}.pptx")

    deck_settings = export_settings_for_deck(photos, settings)
    with tempfile.TemporaryDirectory(prefix=".exif-banner-pptx-", dir=str(output_dir)) as temp_root:
        rendered_images = render_pptx_images_parallel(photos, deck_settings, temp_dir=Path(temp_root))
        write_pptx(output_file, rendered_images, int(deck_settings["slideWidth"]), int(deck_settings["slideHeight"]))
    return {"count": len(photos), "outputFile": str(output_file)}


def start_export_job(payload: dict[str, Any]) -> dict[str, Any]:
    kind = clean_text(payload.get("kind"))
    if kind not in {"images", "pptx"}:
        raise ValueError("Unsupported export type.")

    album = get_album(clean_text(payload.get("albumId")))
    photos = selected_photos(album, payload.get("selection"))
    if not photos:
        raise ValueError("No photos selected for export.")

    output_text = clean_text(payload.get("outputDir"))
    if not output_text:
        raise ValueError("Output folder is required.")
    output_dir = Path(output_text)
    output_dir.mkdir(parents=True, exist_ok=True)

    total = len(photos) if kind == "images" else len(photos) + 1
    job_id = uuid.uuid4().hex
    with EXPORT_LOCK:
        EXPORT_JOBS[job_id] = {
            "id": job_id,
            "kind": kind,
            "status": "running",
            "done": 0,
            "total": total,
            "progress": 0,
            "message": "准备导出",
            "result": None,
            "error": None,
            "cancelRequested": False,
            "created": time.time(),
        }

    worker_payload = dict(payload)
    worker_payload["outputDir"] = str(output_dir)
    Thread(target=run_export_job, args=(job_id, kind, worker_payload), daemon=True).start()
    return {"jobId": job_id, "total": total}


def run_export_job(job_id: str, kind: str, payload: dict[str, Any]) -> None:
    try:
        album = get_album(clean_text(payload.get("albumId")))
        settings = merged_settings(payload.get("settings"))
        photos = selected_photos(album, payload.get("selection"))
        output_dir = Path(clean_text(payload.get("outputDir")))
        if kind == "images":
            result = export_images_with_progress(photos, output_dir, settings, job_id)
        else:
            result = export_pptx_with_progress(album, photos, output_dir, settings, job_id)
        ensure_export_not_cancelled(job_id)
        update_export_job(job_id, status="done", progress=1, message="导出完成", result=result)
    except ExportCancelled:
        update_export_job(job_id, status="canceled", message="导出已取消", error=None)
    except Exception as exc:
        update_export_job(job_id, status="error", message="导出失败", error=str(exc))


def export_images_with_progress(
    photos: list[dict[str, Any]],
    output_dir: Path,
    settings: dict[str, Any],
    job_id: str,
) -> dict[str, Any]:
    output_dir.mkdir(parents=True, exist_ok=True)
    total = max(1, len(photos))
    image_ext = export_extension(settings)
    targets = unique_export_paths(photos, output_dir, image_ext, None)
    written = export_photo_images_parallel(photos, targets, settings, job_id, total, "导出图片")
    return {"count": len(written), "outputDir": str(output_dir), "files": written[:12]}


def export_pptx_with_progress(
    album: dict[str, Any],
    photos: list[dict[str, Any]],
    output_dir: Path,
    settings: dict[str, Any],
    job_id: str,
) -> dict[str, Any]:
    output_dir.mkdir(parents=True, exist_ok=True)
    root = Path(album["root"])
    output_file = unique_path(output_dir / f"{root.name}_banner_{timestamp()}.pptx")
    total = max(1, len(photos) + 1)
    deck_settings = export_settings_for_deck(photos, settings)
    with tempfile.TemporaryDirectory(prefix=".exif-banner-pptx-", dir=str(output_dir)) as temp_root:
        rendered_images = render_pptx_images_parallel(
            photos,
            deck_settings,
            job_id,
            total,
            "渲染幻灯片",
            progress_total=len(photos),
            temp_dir=Path(temp_root),
        )
        ensure_export_not_cancelled(job_id)
        update_export_job(job_id, done=len(photos), total=total, progress=len(photos) / total, message="写入 PPTX")
        write_pptx(output_file, rendered_images, int(deck_settings["slideWidth"]), int(deck_settings["slideHeight"]))
    update_export_job(job_id, done=total, total=total, progress=1, message="写入 PPTX")
    return {"count": len(photos), "outputFile": str(output_file)}


def export_photo_images_parallel(
    photos: list[dict[str, Any]],
    targets: list[Path],
    settings: dict[str, Any],
    job_id: str | None = None,
    total: int | None = None,
    message_prefix: str = "导出图片",
    progress_total: int | None = None,
) -> list[str]:
    if not photos:
        return []
    workers = export_worker_count(settings, len(photos))
    total_units = total or len(photos)
    display_total = progress_total or len(photos)
    if job_id:
        update_export_job(
            job_id,
            done=0,
            total=total_units,
            progress=0,
            message=f"{message_prefix} 0/{display_total}",
        )
    ensure_export_not_cancelled(job_id)

    ordered: list[str | None] = [None] * len(photos)
    completed = 0
    with ThreadPoolExecutor(max_workers=workers) as executor:
        future_map = {
            executor.submit(render_and_save_export_photo, photo, target, settings): index
            for index, (photo, target) in enumerate(zip(photos, targets))
        }
        for future in as_completed(future_map):
            if export_cancel_requested(job_id):
                cancel_pending_futures(future_map)
                raise ExportCancelled()
            index = future_map[future]
            ordered[index] = future.result()
            completed += 1
            if job_id:
                update_export_job(
                    job_id,
                    done=completed,
                    total=total_units,
                    progress=completed / total_units,
                    message=f"{message_prefix} {completed}/{display_total}",
                )
            ensure_export_not_cancelled(job_id)

    return [path for path in ordered if path is not None]


def render_pptx_images_parallel(
    photos: list[dict[str, Any]],
    settings: dict[str, Any],
    job_id: str | None = None,
    total: int | None = None,
    message_prefix: str = "渲染幻灯片",
    progress_total: int | None = None,
    temp_dir: Path | None = None,
) -> list[dict[str, Any]]:
    if not photos:
        return []
    workers = export_worker_count(settings, len(photos))
    total_units = total or len(photos)
    display_total = progress_total or len(photos)
    if job_id:
        update_export_job(
            job_id,
            done=0,
            total=total_units,
            progress=0,
            message=f"{message_prefix} 0/{display_total}",
        )
    ensure_export_not_cancelled(job_id)

    ordered: list[dict[str, Any] | None] = [None] * len(photos)
    completed = 0
    image_ext = export_extension(settings)
    if temp_dir is not None:
        temp_dir.mkdir(parents=True, exist_ok=True)
    with ThreadPoolExecutor(max_workers=workers) as executor:
        if temp_dir is None:
            future_map = {
                executor.submit(render_export_photo_bytes, photo, settings): index
                for index, photo in enumerate(photos)
            }
        else:
            future_map = {
                executor.submit(render_export_photo_file, photo, temp_dir / f"slide_{index + 1:05d}.{image_ext}", settings): index
                for index, photo in enumerate(photos)
            }
        for future in as_completed(future_map):
            if export_cancel_requested(job_id):
                cancel_pending_futures(future_map)
                raise ExportCancelled()
            index = future_map[future]
            ordered[index] = future.result()
            completed += 1
            if job_id:
                update_export_job(
                    job_id,
                    done=completed,
                    total=total_units,
                    progress=completed / total_units,
                    message=f"{message_prefix} {completed}/{display_total}",
                )
            ensure_export_not_cancelled(job_id)

    return [item for item in ordered if item is not None]


def cancel_pending_futures(future_map: dict[Any, int]) -> None:
    for future in future_map:
        future.cancel()


def render_and_save_export_photo(photo: dict[str, Any], target: Path, settings: dict[str, Any]) -> str:
    render_settings = settings if settings.get("_fixedSlideSize") else export_settings_for_photo(photo, settings)
    image = render_composite(photo, render_settings)
    save_export_image(image, target, render_settings)
    return str(target)


def render_export_photo_bytes(photo: dict[str, Any], settings: dict[str, Any]) -> dict[str, Any]:
    render_settings = settings if settings.get("_fixedSlideSize") else export_settings_for_photo(photo, settings)
    image = render_composite(photo, render_settings)
    return export_image_bytes(image, render_settings, for_pptx=True)


def render_export_photo_file(photo: dict[str, Any], target: Path, settings: dict[str, Any]) -> dict[str, Any]:
    render_settings = settings if settings.get("_fixedSlideSize") else export_settings_for_photo(photo, settings)
    image = render_composite(photo, render_settings)
    save_pptx_image(image, target, render_settings)
    return {
        "path": str(target),
        "ext": export_extension(render_settings),
        "mime": export_mime(render_settings),
    }


def export_worker_count(settings: dict[str, Any], total: int) -> int:
    cpu_count = os.cpu_count() or 2
    workers = max(1, cpu_count - 1)
    return min(workers, max(1, total))


def unique_export_paths(
    photos: list[dict[str, Any]],
    output_dir: Path,
    image_ext: str,
    source_root: Path | None,
) -> list[Path]:
    reserved: set[str] = set()
    paths: list[Path] = []
    for photo in photos:
        photo_path = Path(clean_text(photo.get("path")) or photo.get("name", "photo"))
        folder = output_dir / relative_photo_folder(photo_path, source_root)
        base = folder / f"{Path(photo['name']).stem}_banner.{image_ext}"
        target = unique_path_reserved(base, reserved)
        reserved.add(str(target).casefold())
        paths.append(target)
    return paths


def relative_photo_folder(photo_path: Path, source_root: Path | None) -> Path:
    if not source_root:
        return Path()
    try:
        relative = photo_path.resolve().relative_to(source_root.resolve())
    except Exception:
        return Path()
    parent = relative.parent
    if str(parent) == ".":
        return Path()
    return safe_relative_folder(parent)


def safe_relative_folder(folder: Path) -> Path:
    parts = [part for part in folder.parts if part not in {"", ".", ".."}]
    return Path(*parts) if parts else Path()


def unique_path_reserved(path: Path, reserved: set[str]) -> Path:
    if not path.exists() and str(path).casefold() not in reserved:
        return path
    stem = path.stem
    suffix = path.suffix
    for index in range(2, 10000):
        candidate = path.with_name(f"{stem}_{index}{suffix}")
        if not candidate.exists() and str(candidate).casefold() not in reserved:
            return candidate
    candidate = path.with_name(f"{stem}_{uuid.uuid4().hex}{suffix}")
    while str(candidate).casefold() in reserved:
        candidate = path.with_name(f"{stem}_{uuid.uuid4().hex}{suffix}")
    return candidate


def save_png(image: Image.Image, target: Path) -> None:
    if image.mode not in {"RGB", "RGBA"}:
        image = image.convert("RGB")
    image.save(target, "PNG", compress_level=4)


def save_export_image(image: Image.Image, target: Path, settings: dict[str, Any]) -> None:
    target.parent.mkdir(parents=True, exist_ok=True)
    if export_format(settings) == "png":
        save_png(image, target)
    else:
        save_jpeg(image, target, settings)


def save_jpeg(image: Image.Image, target: Path, settings: dict[str, Any]) -> None:
    if image.mode != "RGB":
        image = image.convert("RGB")
    image.save(target, "JPEG", **jpeg_save_options(settings))


def save_pptx_image(image: Image.Image, target: Path, settings: dict[str, Any]) -> None:
    target.parent.mkdir(parents=True, exist_ok=True)
    if export_format(settings) == "png":
        save_png(image, target)
        return
    if image.mode != "RGB":
        image = image.convert("RGB")
    image.save(target, "JPEG", **jpeg_save_options(settings, for_pptx=True))


def export_image_bytes(image: Image.Image, settings: dict[str, Any], for_pptx: bool = False) -> dict[str, Any]:
    output = io.BytesIO()
    if export_format(settings) == "png":
        if image.mode not in {"RGB", "RGBA"}:
            image = image.convert("RGB")
        image.save(output, "PNG", compress_level=4)
        return {"data": output.getvalue(), "ext": "png", "mime": "image/png"}

    if image.mode != "RGB":
        image = image.convert("RGB")
    image.save(output, "JPEG", **jpeg_save_options(settings, for_pptx=for_pptx))
    return {"data": output.getvalue(), "ext": "jpg", "mime": "image/jpeg"}


def jpeg_save_options(settings: dict[str, Any], for_pptx: bool = False) -> dict[str, Any]:
    quality = clamp_int(settings.get("quality"), 60, 100, 92)
    subsampling = 0 if quality >= 96 else 2
    options: dict[str, Any] = {"quality": quality, "subsampling": subsampling}
    if not for_pptx:
        options["progressive"] = True
    return options


def export_format(settings: dict[str, Any]) -> str:
    value = clean_text(settings.get("exportFormat")).lower()
    if value == "png":
        return "png"
    return "jpeg"


def export_mime(settings: dict[str, Any]) -> str:
    return "image/png" if export_format(settings) == "png" else "image/jpeg"


def export_extension(settings: dict[str, Any]) -> str:
    return "png" if export_format(settings) == "png" else "jpg"


def unique_path(path: Path) -> Path:
    if not path.exists():
        return path
    stem = path.stem
    suffix = path.suffix
    for index in range(2, 10000):
        candidate = path.with_name(f"{stem}_{index}{suffix}")
        if not candidate.exists():
            return candidate
    return path.with_name(f"{stem}_{uuid.uuid4().hex}{suffix}")


def update_export_job(job_id: str, **updates: Any) -> None:
    with EXPORT_LOCK:
        job = EXPORT_JOBS.get(job_id)
        if not job:
            return
        job.update(updates)


def request_export_cancel(job_id: str) -> dict[str, Any]:
    with EXPORT_LOCK:
        job = EXPORT_JOBS.get(job_id)
        if not job:
            raise ValueError("Export job not found.")
        if job.get("status") == "running":
            job["cancelRequested"] = True
            job["message"] = "正在取消"
        return dict(job)


def export_cancel_requested(job_id: str | None) -> bool:
    if not job_id:
        return False
    with EXPORT_LOCK:
        job = EXPORT_JOBS.get(job_id)
        return bool(job and job.get("cancelRequested"))


def ensure_export_not_cancelled(job_id: str | None) -> None:
    if export_cancel_requested(job_id):
        raise ExportCancelled()


def get_export_job(job_id: str) -> dict[str, Any]:
    with EXPORT_LOCK:
        job = EXPORT_JOBS.get(job_id)
        if not job:
            raise ValueError("Export job not found.")
        return dict(job)


def get_album(album_id: str) -> dict[str, Any]:
    album = SESSIONS.get(album_id)
    if not album:
        raise ValueError("Album session expired. Scan the folder again.")
    return album


def timestamp() -> str:
    return datetime.now().strftime("%Y%m%d_%H%M%S")


EMU_PER_INCH = 914400


def px_to_emu(px: int, dpi: int = 144) -> int:
    return int(px / dpi * EMU_PER_INCH)


def write_pptx(output_file: Path, rendered_images: list[dict[str, Any]], width_px: int, height_px: int) -> None:
    if Presentation is None:
        raise RuntimeError("缺少 python-pptx 依赖。请运行 pip install -r requirements.txt 后再导出 PPTX。")

    presentation = Presentation()
    presentation.slide_width = px_to_emu(width_px)
    presentation.slide_height = px_to_emu(height_px)
    blank_layout = presentation.slide_layouts[6]

    for rendered in rendered_images:
        slide = presentation.slides.add_slide(blank_layout)
        image_source = str(rendered["path"]) if rendered.get("path") else io.BytesIO(rendered["data"])
        slide.shapes.add_picture(
            image_source,
            0,
            0,
            width=presentation.slide_width,
            height=presentation.slide_height,
        )

    output_file.parent.mkdir(parents=True, exist_ok=True)
    presentation.save(output_file)


def content_types(slide_count: int, image_ext: str, image_mime: str) -> str:
    slide_overrides = "\n".join(
        f'<Override PartName="/ppt/slides/slide{index}.xml" '
        'ContentType="application/vnd.openxmlformats-officedocument.presentationml.slide+xml"/>'
        for index in range(1, slide_count + 1)
    )
    return f'''<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<Types xmlns="http://schemas.openxmlformats.org/package/2006/content-types">
  <Default Extension="rels" ContentType="application/vnd.openxmlformats-package.relationships+xml"/>
  <Default Extension="xml" ContentType="application/xml"/>
  <Default Extension="{image_ext}" ContentType="{image_mime}"/>
  <Override PartName="/docProps/app.xml" ContentType="application/vnd.openxmlformats-officedocument.extended-properties+xml"/>
  <Override PartName="/docProps/core.xml" ContentType="application/vnd.openxmlformats-package.core-properties+xml"/>
  <Override PartName="/ppt/presentation.xml" ContentType="application/vnd.openxmlformats-officedocument.presentationml.presentation.main+xml"/>
  <Override PartName="/ppt/slideMasters/slideMaster1.xml" ContentType="application/vnd.openxmlformats-officedocument.presentationml.slideMaster+xml"/>
  <Override PartName="/ppt/slideLayouts/slideLayout1.xml" ContentType="application/vnd.openxmlformats-officedocument.presentationml.slideLayout+xml"/>
  <Override PartName="/ppt/theme/theme1.xml" ContentType="application/vnd.openxmlformats-officedocument.theme+xml"/>
  {slide_overrides}
</Types>'''


def package_rels() -> str:
    return '''<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">
  <Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/officeDocument" Target="ppt/presentation.xml"/>
  <Relationship Id="rId2" Type="http://schemas.openxmlformats.org/package/2006/relationships/metadata/core-properties" Target="docProps/core.xml"/>
  <Relationship Id="rId3" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/extended-properties" Target="docProps/app.xml"/>
</Relationships>'''


def app_xml(slide_count: int) -> str:
    return f'''<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<Properties xmlns="http://schemas.openxmlformats.org/officeDocument/2006/extended-properties"
  xmlns:vt="http://schemas.openxmlformats.org/officeDocument/2006/docPropsVTypes">
  <Application>EXIF-Banner</Application>
  <PresentationFormat>Widescreen</PresentationFormat>
  <Slides>{slide_count}</Slides>
  <Company></Company>
</Properties>'''


def core_xml() -> str:
    now = datetime.utcnow().strftime("%Y-%m-%dT%H:%M:%SZ")
    return f'''<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<cp:coreProperties xmlns:cp="http://schemas.openxmlformats.org/package/2006/metadata/core-properties"
  xmlns:dc="http://purl.org/dc/elements/1.1/"
  xmlns:dcterms="http://purl.org/dc/terms/"
  xmlns:dcmitype="http://purl.org/dc/dcmitype/"
  xmlns:xsi="http://www.w3.org/2001/XMLSchema-instance">
  <dc:title>EXIF-Banner Export</dc:title>
  <dc:creator>EXIF-Banner</dc:creator>
  <cp:lastModifiedBy>EXIF-Banner</cp:lastModifiedBy>
  <dcterms:created xsi:type="dcterms:W3CDTF">{now}</dcterms:created>
  <dcterms:modified xsi:type="dcterms:W3CDTF">{now}</dcterms:modified>
</cp:coreProperties>'''


def presentation_xml(slide_count: int, slide_cx: int, slide_cy: int) -> str:
    slide_ids = "\n".join(
        f'<p:sldId id="{255 + index}" r:id="rId{index + 1}"/>'
        for index in range(1, slide_count + 1)
    )
    return f'''<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<p:presentation xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"
  xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships"
  xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main">
  <p:sldMasterIdLst><p:sldMasterId id="2147483648" r:id="rId1"/></p:sldMasterIdLst>
  <p:sldIdLst>{slide_ids}</p:sldIdLst>
  <p:sldSz cx="{slide_cx}" cy="{slide_cy}" type="screen16x9"/>
  <p:notesSz cx="6858000" cy="9144000"/>
</p:presentation>'''


def presentation_rels(slide_count: int) -> str:
    rels = [
        '<Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/slideMaster" Target="slideMasters/slideMaster1.xml"/>'
    ]
    for index in range(1, slide_count + 1):
        rels.append(
            f'<Relationship Id="rId{index + 1}" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/slide" Target="slides/slide{index}.xml"/>'
        )
    return '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>\n<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">\n  ' + "\n  ".join(rels) + "\n</Relationships>"


def slide_xml(index: int, slide_cx: int, slide_cy: int) -> str:
    return f'''<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<p:sld xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"
  xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships"
  xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main">
  <p:cSld>
    <p:spTree>
      <p:nvGrpSpPr>
        <p:cNvPr id="1" name=""/>
        <p:cNvGrpSpPr/>
        <p:nvPr/>
      </p:nvGrpSpPr>
      <p:grpSpPr>
        <a:xfrm><a:off x="0" y="0"/><a:ext cx="0" cy="0"/><a:chOff x="0" y="0"/><a:chExt cx="0" cy="0"/></a:xfrm>
      </p:grpSpPr>
      <p:pic>
        <p:nvPicPr>
          <p:cNvPr id="2" name="Banner {index}"/>
          <p:cNvPicPr><a:picLocks noChangeAspect="1"/></p:cNvPicPr>
          <p:nvPr/>
        </p:nvPicPr>
        <p:blipFill>
          <a:blip r:embed="rId1"/>
          <a:stretch><a:fillRect/></a:stretch>
        </p:blipFill>
        <p:spPr>
          <a:xfrm><a:off x="0" y="0"/><a:ext cx="{slide_cx}" cy="{slide_cy}"/></a:xfrm>
          <a:prstGeom prst="rect"><a:avLst/></a:prstGeom>
        </p:spPr>
      </p:pic>
    </p:spTree>
  </p:cSld>
  <p:clrMapOvr><a:masterClrMapping/></p:clrMapOvr>
</p:sld>'''


def slide_rels(index: int, image_ext: str) -> str:
    return f'''<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">
  <Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/image" Target="../media/image{index}.{image_ext}"/>
  <Relationship Id="rId2" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/slideLayout" Target="../slideLayouts/slideLayout1.xml"/>
</Relationships>'''


def slide_master_xml() -> str:
    return '''<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<p:sldMaster xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"
  xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships"
  xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main">
  <p:cSld><p:spTree><p:nvGrpSpPr><p:cNvPr id="1" name=""/><p:cNvGrpSpPr/><p:nvPr/></p:nvGrpSpPr><p:grpSpPr><a:xfrm><a:off x="0" y="0"/><a:ext cx="0" cy="0"/><a:chOff x="0" y="0"/><a:chExt cx="0" cy="0"/></a:xfrm></p:grpSpPr></p:spTree></p:cSld>
  <p:clrMap bg1="lt1" tx1="dk1" bg2="lt2" tx2="dk2" accent1="accent1" accent2="accent2" accent3="accent3" accent4="accent4" accent5="accent5" accent6="accent6" hlink="hlink" folHlink="folHlink"/>
  <p:sldLayoutIdLst><p:sldLayoutId id="2147483649" r:id="rId1"/></p:sldLayoutIdLst>
  <p:txStyles><p:titleStyle/><p:bodyStyle/><p:otherStyle/></p:txStyles>
</p:sldMaster>'''


def slide_master_rels() -> str:
    return '''<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">
  <Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/slideLayout" Target="../slideLayouts/slideLayout1.xml"/>
  <Relationship Id="rId2" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/theme" Target="../theme/theme1.xml"/>
</Relationships>'''


def slide_layout_xml() -> str:
    return '''<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<p:sldLayout xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"
  xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships"
  xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" type="blank" preserve="1">
  <p:cSld name="Blank"><p:spTree><p:nvGrpSpPr><p:cNvPr id="1" name=""/><p:cNvGrpSpPr/><p:nvPr/></p:nvGrpSpPr><p:grpSpPr><a:xfrm><a:off x="0" y="0"/><a:ext cx="0" cy="0"/><a:chOff x="0" y="0"/><a:chExt cx="0" cy="0"/></a:xfrm></p:grpSpPr></p:spTree></p:cSld>
  <p:clrMapOvr><a:masterClrMapping/></p:clrMapOvr>
</p:sldLayout>'''


def slide_layout_rels() -> str:
    return '''<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">
  <Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/slideMaster" Target="../slideMasters/slideMaster1.xml"/>
</Relationships>'''


def theme_xml() -> str:
    return '''<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<a:theme xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" name="EXIF-Banner">
  <a:themeElements>
    <a:clrScheme name="Office">
      <a:dk1><a:sysClr val="windowText" lastClr="000000"/></a:dk1>
      <a:lt1><a:sysClr val="window" lastClr="FFFFFF"/></a:lt1>
      <a:dk2><a:srgbClr val="1F1F1F"/></a:dk2>
      <a:lt2><a:srgbClr val="F6F4EF"/></a:lt2>
      <a:accent1><a:srgbClr val="4B6B5B"/></a:accent1>
      <a:accent2><a:srgbClr val="B58B5B"/></a:accent2>
      <a:accent3><a:srgbClr val="6C7480"/></a:accent3>
      <a:accent4><a:srgbClr val="8B4A47"/></a:accent4>
      <a:accent5><a:srgbClr val="57708A"/></a:accent5>
      <a:accent6><a:srgbClr val="7A6B53"/></a:accent6>
      <a:hlink><a:srgbClr val="0563C1"/></a:hlink>
      <a:folHlink><a:srgbClr val="954F72"/></a:folHlink>
    </a:clrScheme>
    <a:fontScheme name="Office"><a:majorFont><a:latin typeface="Segoe UI"/></a:majorFont><a:minorFont><a:latin typeface="Segoe UI"/></a:minorFont></a:fontScheme>
    <a:fmtScheme name="Office"><a:fillStyleLst><a:solidFill><a:schemeClr val="phClr"/></a:solidFill></a:fillStyleLst><a:lnStyleLst><a:ln w="6350"><a:solidFill><a:schemeClr val="phClr"/></a:solidFill></a:ln></a:lnStyleLst><a:effectStyleLst><a:effectStyle/></a:effectStyleLst><a:bgFillStyleLst><a:solidFill><a:schemeClr val="phClr"/></a:solidFill></a:bgFillStyleLst></a:fmtScheme>
  </a:themeElements>
</a:theme>'''


class Handler(SimpleHTTPRequestHandler):
    server_version = "EXIFBanner/0.1"

    def do_GET(self) -> None:
        parsed = urlparse(self.path)
        if parsed.path == "/":
            return self.serve_file(STATIC_DIR / "index.html")
        if parsed.path.startswith("/static/"):
            target = (STATIC_DIR / parsed.path.removeprefix("/static/")).resolve()
            if not str(target).startswith(str(STATIC_DIR.resolve())):
                return self.send_error(404)
            return self.serve_file(target)
        if parsed.path.startswith("/logos/"):
            target = (LOGO_DIR / parsed.path.removeprefix("/logos/")).resolve()
            if not str(target).startswith(str(LOGO_DIR.resolve())):
                return self.send_error(404)
            return self.serve_file(target)
        if parsed.path == "/api/photo":
            return self.serve_photo(parsed.query)
        if parsed.path == "/api/logo":
            return self.serve_logo(parsed.query)
        if parsed.path == "/api/status":
            return self.send_json({"ok": True, "settings": DEFAULT_SETTINGS})
        self.send_error(404)

    def do_POST(self) -> None:
        parsed = urlparse(self.path)
        try:
            payload = self.read_json()
            if parsed.path == "/api/pick-folder":
                return self.send_json({"folder": pick_folder()})
            if parsed.path == "/api/pick-logo":
                return self.send_json({"file": pick_logo()})
            if parsed.path == "/api/pick-output-dir":
                return self.send_json({"folder": pick_output_dir()})
            if parsed.path == "/api/scan":
                folder = Path(clean_text(payload.get("folder"))).expanduser()
                recursive = bool(payload.get("recursive"))
                return self.send_json(scan_album(folder, recursive, clean_text(payload.get("sortMode"))))
            if parsed.path == "/api/export/start":
                return self.send_json(start_export_job(payload))
            if parsed.path == "/api/export/status":
                return self.send_json(get_export_job(clean_text(payload.get("jobId"))))
            if parsed.path == "/api/export/cancel":
                return self.send_json(request_export_cancel(clean_text(payload.get("jobId"))))
            if parsed.path == "/api/banner-layout":
                return self.send_json(self.banner_layout(payload))
            if parsed.path == "/api/preview":
                return self.send_preview(payload)
            if parsed.path == "/api/export/images":
                return self.send_json(export_images(clean_text(payload.get("albumId")), payload))
            if parsed.path == "/api/export/pptx":
                return self.send_json(export_pptx(clean_text(payload.get("albumId")), payload))
            self.send_error(404)
        except Exception as exc:
            self.send_json({"error": str(exc)}, status=400)

    def read_json(self) -> dict[str, Any]:
        length = int(self.headers.get("Content-Length", "0") or "0")
        if not length:
            return {}
        body = self.rfile.read(length).decode("utf-8", errors="replace")
        return json.loads(body or "{}")

    def send_json(self, payload: dict[str, Any], status: int = 200) -> None:
        data = json.dumps(payload, ensure_ascii=False).encode("utf-8")
        self.send_response(status)
        self.send_header("Content-Type", "application/json; charset=utf-8")
        self.send_header("Content-Length", str(len(data)))
        self.end_headers()
        self.wfile.write(data)

    def banner_layout(self, payload: dict[str, Any]) -> dict[str, Any]:
        album_id = clean_text(payload.get("albumId"))
        album = get_album(album_id)
        index = int(payload.get("index", 0))
        photos = album["photos"]
        if index < 0 or index >= len(photos):
            raise ValueError("Banner layout photo index is out of range.")
        return banner_edit_layout(photos[index], payload.get("settings"))

    def send_preview(self, payload: dict[str, Any]) -> None:
        album_id = clean_text(payload.get("albumId"))
        album = get_album(album_id)
        index = int(payload.get("index", 0))
        photos = album["photos"]
        if index < 0 or index >= len(photos):
            raise ValueError("Preview photo index is out of range.")
        settings = merged_settings(payload.get("settings"))
        max_size = clamp_int(payload.get("maxSize"), 600, 2400, 1600)
        cache_key = preview_cache_key(album_id, index, photos[index], settings, max_size)
        data = get_cached_preview(cache_key)
        if data is None:
            image = render_preview_composite(photos[index], settings, max_size)
            image.thumbnail((max_size, max_size), RESAMPLE)
            output = io.BytesIO()
            image.save(output, "JPEG", quality=88, subsampling=2)
            data = output.getvalue()
            remember_preview(cache_key, data)
        self.send_response(200)
        self.send_header("Content-Type", "image/jpeg")
        self.send_header("Cache-Control", "no-store")
        self.send_header("Content-Length", str(len(data)))
        self.end_headers()
        self.wfile.write(data)

    def serve_file(self, path: Path) -> None:
        if not path.exists() or not path.is_file():
            return self.send_error(404)
        data = path.read_bytes()
        content_type = mimetypes.guess_type(str(path))[0] or "application/octet-stream"
        self.send_response(200)
        self.send_header("Content-Type", content_type)
        self.send_header("Cache-Control", "no-store")
        self.send_header("Content-Length", str(len(data)))
        self.end_headers()
        self.wfile.write(data)

    def serve_photo(self, query: str) -> None:
        params = parse_qs(query)
        album = get_album(params.get("album", [""])[0])
        index = int(params.get("index", ["0"])[0])
        max_size = clamp_int(params.get("max", ["1600"])[0], 100, 3000, 1600)
        photo = album["photos"][index]
        image = load_oriented_rgb(photo["path"], (max_size, max_size))
        image.thumbnail((max_size, max_size), RESAMPLE)
        output = io.BytesIO()
        image.save(output, "JPEG", quality=82, subsampling=2)
        data = output.getvalue()
        self.send_response(200)
        self.send_header("Content-Type", "image/jpeg")
        self.send_header("Cache-Control", "no-store")
        self.send_header("Content-Length", str(len(data)))
        self.end_headers()
        self.wfile.write(data)

    def serve_logo(self, query: str) -> None:
        params = parse_qs(query)
        image = load_logo_image(params.get("path", [""])[0])
        if image is None:
            return self.send_error(404)
        image.thumbnail((800, 400), RESAMPLE)
        output = io.BytesIO()
        image.save(output, "PNG")
        data = output.getvalue()
        self.send_response(200)
        self.send_header("Content-Type", "image/png")
        self.send_header("Cache-Control", "no-store")
        self.send_header("Content-Length", str(len(data)))
        self.end_headers()
        self.wfile.write(data)


def pick_folder() -> str:
    try:
        import tkinter as tk
        from tkinter import filedialog

        root = tk.Tk()
        root.withdraw()
        root.attributes("-topmost", True)
        folder = filedialog.askdirectory(title="Select photo folder")
        root.destroy()
        return folder or ""
    except Exception as exc:
        raise RuntimeError(f"Folder picker is unavailable: {exc}") from exc


def pick_logo() -> str:
    try:
        import tkinter as tk
        from tkinter import filedialog

        root = tk.Tk()
        root.withdraw()
        root.attributes("-topmost", True)
        file_name = filedialog.askopenfilename(
            title="Select logo image",
            filetypes=[
                ("Image files", "*.png;*.jpg;*.jpeg;*.webp;*.bmp"),
                ("All files", "*.*"),
            ],
        )
        root.destroy()
        return file_name or ""
    except Exception as exc:
        raise RuntimeError(f"Logo picker is unavailable: {exc}") from exc


def pick_output_dir() -> str:
    try:
        import tkinter as tk
        from tkinter import filedialog

        root = tk.Tk()
        root.withdraw()
        root.attributes("-topmost", True)
        folder = filedialog.askdirectory(title="Select export folder")
        root.destroy()
        return folder or ""
    except Exception as exc:
        raise RuntimeError(f"Output folder picker is unavailable: {exc}") from exc


def main() -> None:
    parser = argparse.ArgumentParser(description="Run the EXIF-Banner local web app.")
    parser.add_argument("--host", default="127.0.0.1")
    parser.add_argument("--port", type=int, default=8765)
    parser.add_argument("--no-browser", action="store_true")
    args = parser.parse_args()

    server = ThreadingHTTPServer((args.host, args.port), Handler)
    url = f"http://{args.host}:{args.port}/"
    print(f"EXIF-Banner is running at {url}")
    if not args.no_browser:
        Timer(0.8, lambda: webbrowser.open(url)).start()
    try:
        server.serve_forever()
    except KeyboardInterrupt:
        pass
    finally:
        server.server_close()


if __name__ == "__main__":
    main()
